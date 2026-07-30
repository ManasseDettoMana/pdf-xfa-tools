"""Document conversions: office formats, text, Markdown, HTML, spreadsheets.

Office-to-PDF is the one family with no honest pure-Python answer.  Three
back-ends are tried in descending order of fidelity:

1. **LibreOffice** headless - renders with the real layout engine.
2. **Microsoft Word** via COM - Windows only, needs Word and ``pywin32``.
3. **ReportLab** - a plain-text reflow of the document's content.

Only the first two preserve layout.  The ReportLab path is a readable fallback,
not a faithful rendering, and every result it produces is flagged as approximate
so the user knows what they got.
"""

from __future__ import annotations

import shutil
import tempfile
from pathlib import Path
from typing import Any

from . import engines
from .errors import ConversionError, EngineMissingError
from .job import NULL_CONTEXT, JobContext

OFFICE_EXTS = ("docx", "doc", "odt", "rtf", "xlsx", "xls", "ods", "pptx", "ppt", "odp")
TEXT_EXTS = ("txt", "md", "markdown", "html", "htm", "csv")

#: Page geometry for the ReportLab fallback, in points.
FALLBACK_MARGIN = 54.0


# ---------------------------------------------------------------------------
# Office -> PDF
# ---------------------------------------------------------------------------


def _libreoffice_convert(src: Path, dst: Path, target: str, ctx: JobContext) -> Path:
    """Convert via headless LibreOffice into ``dst``."""
    soffice = engines.find_libreoffice()
    if not soffice:
        raise EngineMissingError("LibreOffice", "la conversione dei documenti", engines.LIBREOFFICE_HINT)

    ctx.progress(0, 1, f"{src.name}: conversione con LibreOffice")

    # A private profile directory lets us run even while the user has
    # LibreOffice open, which otherwise makes the headless call exit silently.
    with tempfile.TemporaryDirectory(prefix="xfatools-lo-") as work_dir:
        profile = Path(work_dir) / "profile"
        out_dir = Path(work_dir) / "out"
        out_dir.mkdir()

        result = engines.run_hidden(
            [
                soffice,
                f"-env:UserInstallation={profile.as_uri()}",
                "--headless",
                "--norestore",
                "--convert-to",
                target,
                "--outdir",
                str(out_dir),
                str(src),
            ],
            timeout=300,
        )

        produced = list(out_dir.glob(f"*.{target.split(':')[0]}"))
        if not produced:
            detail = (result.stderr or result.stdout or "").strip()[:300]
            raise ConversionError(
                f"LibreOffice non ha prodotto un file per '{src.name}'."
                + (f" Dettagli: {detail}" if detail else "")
            )

        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(produced[0]), str(dst))

    return dst


def _word_com_convert(src: Path, dst: Path, ctx: JobContext) -> Path:
    """Convert a Word document to PDF by driving an installed Word."""
    if not engines.has_word_com():
        raise EngineMissingError("Microsoft Word", "la conversione DOCX", engines.WORD_HINT)

    import pythoncom
    import win32com.client

    ctx.progress(0, 1, f"{src.name}: conversione con Microsoft Word")

    wd_format_pdf = 17
    dst.parent.mkdir(parents=True, exist_ok=True)

    # The worker runs on a non-main thread, so COM has to be initialised here.
    pythoncom.CoInitialize()
    word = None
    document = None
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        document = word.Documents.Open(str(src.resolve()), ReadOnly=True)
        document.SaveAs(str(dst.resolve()), FileFormat=wd_format_pdf)
    except Exception as exc:
        raise ConversionError(f"Word non e' riuscito a convertire '{src.name}': {exc}") from exc
    finally:
        if document is not None:
            document.Close(False)
        if word is not None:
            word.Quit()
        pythoncom.CoUninitialize()

    return dst


def office_to_pdf(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Convert an office document to PDF using the best available back-end."""
    src, dst = Path(src), Path(dst)
    ctx.check_cancelled()

    if engines.find_libreoffice():
        return [_libreoffice_convert(src, dst, "pdf", ctx)]

    if src.suffix.lower() in (".docx", ".doc", ".rtf", ".odt") and engines.has_word_com():
        return [_word_com_convert(src, dst, ctx)]

    # Last resort: extract what text we can and lay it out plainly.
    if src.suffix.lower() in (".docx", ".xlsx", ".pptx"):
        ctx.progress(0, 1, f"{src.name}: nessun motore installato, layout approssimato")
        html = _office_to_html(src)
        _html_to_pdf_fallback(html, dst, title=src.stem)
        return [dst]

    raise EngineMissingError(
        "LibreOffice", f"la conversione di '{src.suffix}'", engines.LIBREOFFICE_HINT
    )


def _office_to_html(src: Path) -> str:
    """Best-effort HTML from an office file, without any external engine."""
    suffix = src.suffix.lower()

    if suffix == ".docx":
        import mammoth

        with src.open("rb") as handle:
            return mammoth.convert_to_html(handle).value

    if suffix == ".xlsx":
        import openpyxl

        workbook = openpyxl.load_workbook(src, data_only=True)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            parts.append(f"<h2>{_escape(sheet.title)}</h2>")
            for row in sheet.iter_rows(values_only=True):
                cells = " | ".join("" if cell is None else str(cell) for cell in row)
                if cells.strip(" |"):
                    parts.append(f"<p>{_escape(cells)}</p>")
        workbook.close()
        return "\n".join(parts)

    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(src))
        parts = []
        for number, slide in enumerate(presentation.slides, start=1):
            parts.append(f"<h2>Slide {number}</h2>")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(f"<p>{_escape(shape.text_frame.text)}</p>")
        return "\n".join(parts)

    raise ConversionError(f"Nessun fallback disponibile per '{src.suffix}'.")


# ---------------------------------------------------------------------------
# Text-like -> PDF
# ---------------------------------------------------------------------------


def text_to_pdf(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Convert txt / md / html / csv to PDF.

    LibreOffice is used when available (it honours real HTML styling); otherwise
    ReportLab lays the content out in a clean single column.
    """
    src, dst = Path(src), Path(dst)
    suffix = src.suffix.lower()
    ctx.check_cancelled()

    if engines.find_libreoffice() and suffix in (".html", ".htm", ".txt", ".csv"):
        return [_libreoffice_convert(src, dst, "pdf", ctx)]

    text = src.read_text(encoding="utf-8", errors="replace")

    if suffix in (".md", ".markdown"):
        html = _markdown_to_html(text)
    elif suffix in (".html", ".htm"):
        html = text
    elif suffix == ".csv":
        html = _csv_to_html(text)
    else:
        html = "\n".join(f"<p>{_escape(line)}</p>" for line in text.splitlines() if line.strip())

    ctx.progress(0, 1, f"{src.name}: generazione PDF")
    _html_to_pdf_fallback(html, dst, title=src.stem)
    ctx.progress(1, 1, f"{src.name}: completato")
    return [dst]


def _markdown_to_html(text: str) -> str:
    try:
        import markdown

        return markdown.markdown(text, extensions=["tables", "fenced_code"])
    except ImportError:
        return "\n".join(f"<p>{_escape(line)}</p>" for line in text.splitlines() if line.strip())


def _csv_to_html(text: str) -> str:
    import csv
    import io

    rows = list(csv.reader(io.StringIO(text)))
    if not rows:
        return "<p></p>"
    parts = ["<table>"]
    for index, row in enumerate(rows):
        tag = "th" if index == 0 else "td"
        cells = "".join(f"<{tag}>{_escape(cell)}</{tag}>" for cell in row)
        parts.append(f"<tr>{cells}</tr>")
    parts.append("</table>")
    return "\n".join(parts)


def _escape(text: str) -> str:
    from xml.sax.saxutils import escape

    return escape(str(text))


def _html_to_pdf_fallback(html: str, dst: Path, title: str = "") -> Path:
    """Lay out HTML as a PDF with ReportLab.

    This understands headings, paragraphs, lists and tables.  It is a reflow, not
    a rendering: CSS, floats and absolute positioning are ignored.  Callers must
    flag the result as approximate.
    """
    from bs4 import BeautifulSoup
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        ListFlowable,
        ListItem,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        TableStyle,
    )
    from reportlab.platypus import Table as RLTable

    soup = BeautifulSoup(html, "html.parser")
    styles = getSampleStyleSheet()
    story: list[Any] = []

    if title:
        story.append(Paragraph(_escape(title), styles["Title"]))
        story.append(Spacer(1, 6 * mm))

    heading_styles = {
        "h1": styles["Heading1"],
        "h2": styles["Heading2"],
        "h3": styles["Heading3"],
        "h4": styles["Heading4"],
        "h5": styles["Heading5"],
        "h6": styles["Heading6"],
    }

    body = soup.body or soup
    for element in body.find_all(
        ["h1", "h2", "h3", "h4", "h5", "h6", "p", "ul", "ol", "table", "pre"], recursive=True
    ):
        if element.find_parent(["table", "ul", "ol"]):
            continue  # already emitted as part of its container

        name = element.name
        if name in heading_styles:
            text = element.get_text(" ", strip=True)
            if text:
                story.append(Paragraph(_escape(text), heading_styles[name]))
        elif name == "pre":
            text = element.get_text("\n", strip=False)
            if text.strip():
                story.append(Paragraph(_escape(text).replace("\n", "<br/>"), styles["Code"]))
        elif name in ("ul", "ol"):
            items = [
                ListItem(Paragraph(_escape(li.get_text(" ", strip=True)), styles["BodyText"]))
                for li in element.find_all("li", recursive=False)
                if li.get_text(strip=True)
            ]
            if items:
                story.append(
                    ListFlowable(items, bulletType="1" if name == "ol" else "bullet")
                )
        elif name == "table":
            data = [
                [_escape(cell.get_text(" ", strip=True)) for cell in row.find_all(["td", "th"])]
                for row in element.find_all("tr")
            ]
            data = [row for row in data if any(cell for cell in row)]
            if data:
                width = max(len(row) for row in data)
                data = [row + [""] * (width - len(row)) for row in data]
                table = RLTable(
                    [[Paragraph(cell, styles["BodyText"]) for cell in row] for row in data],
                    repeatRows=1,
                )
                table.setStyle(
                    TableStyle(
                        [
                            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#B0B7C3")),
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF1F6")),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                            ("LEFTPADDING", (0, 0), (-1, -1), 4),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                        ]
                    )
                )
                story.append(table)
                story.append(Spacer(1, 4 * mm))
        else:
            text = element.get_text(" ", strip=True)
            if text:
                story.append(Paragraph(_escape(text), styles["BodyText"]))
                story.append(Spacer(1, 2 * mm))

    if not story:
        story.append(Paragraph("(documento vuoto)", styles["BodyText"]))

    dst.parent.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(
        str(dst),
        pagesize=A4,
        leftMargin=FALLBACK_MARGIN,
        rightMargin=FALLBACK_MARGIN,
        topMargin=FALLBACK_MARGIN,
        bottomMargin=FALLBACK_MARGIN,
        title=title or dst.stem,
    ).build(story)
    return dst


# ---------------------------------------------------------------------------
# PDF -> text-like
# ---------------------------------------------------------------------------


def pdf_to_text(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Extract the text layer, one page per section."""
    import pdfplumber

    options = options or {}
    src, dst = Path(src), Path(dst)
    max_pages = options.get("max_pages")
    parts: list[str] = []

    with pdfplumber.open(str(src)) as doc:
        total = len(doc.pages)
        limit = min(total, int(max_pages)) if max_pages else total
        for index, page in enumerate(doc.pages[:limit], start=1):
            ctx.check_cancelled()
            ctx.progress(index - 1, limit, f"{src.name}: pagina {index}/{limit}")
            parts.append(page.extract_text() or "")

    text = "\n\n".join(
        f"--- pagina {number} ---\n{content}" for number, content in enumerate(parts, start=1)
    )
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(text, encoding="utf-8")
    return [dst]


def pdf_to_markdown(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Extract text as Markdown, promoting short standalone lines to headings."""
    import pdfplumber

    options = options or {}
    src, dst = Path(src), Path(dst)
    max_pages = options.get("max_pages")
    lines: list[str] = [f"# {src.stem}", ""]

    with pdfplumber.open(str(src)) as doc:
        total = len(doc.pages)
        limit = min(total, int(max_pages)) if max_pages else total
        for index, page in enumerate(doc.pages[:limit], start=1):
            ctx.check_cancelled()
            ctx.progress(index - 1, limit, f"{src.name}: pagina {index}/{limit}")
            lines.append(f"## Pagina {index}")
            lines.append("")
            for raw in (page.extract_text() or "").splitlines():
                stripped = raw.strip()
                if not stripped:
                    lines.append("")
                elif len(stripped) < 60 and stripped.isupper():
                    lines.append(f"### {stripped.title()}")
                else:
                    lines.append(stripped)
            lines.append("")

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text("\n".join(lines), encoding="utf-8")
    return [dst]


def pdf_to_csv(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Extract every detected table into a single CSV."""
    import csv

    import pdfplumber

    options = options or {}
    src, dst = Path(src), Path(dst)
    max_pages = options.get("max_pages")
    rows: list[list[str]] = []

    with pdfplumber.open(str(src)) as doc:
        total = len(doc.pages)
        limit = min(total, int(max_pages)) if max_pages else total
        for index, page in enumerate(doc.pages[:limit], start=1):
            ctx.check_cancelled()
            ctx.progress(index - 1, limit, f"{src.name}: pagina {index}/{limit}")
            for table in page.extract_tables():
                for row in table:
                    rows.append([(cell or "").replace("\n", " ").strip() for cell in row])
                rows.append([])

    if not rows:
        raise ConversionError(
            f"Nessuna tabella rilevata in '{src.name}'.",
            hint="Prova la conversione in testo o in Markdown.",
        )

    dst.parent.mkdir(parents=True, exist_ok=True)
    with dst.open("w", encoding="utf-8-sig", newline="") as handle:
        csv.writer(handle).writerows(rows)
    return [dst]


# ---------------------------------------------------------------------------
# Office -> text-like
# ---------------------------------------------------------------------------


def office_to_text(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Extract plain text or Markdown from an office document."""
    from bs4 import BeautifulSoup

    src, dst = Path(src), Path(dst)
    ctx.progress(0, 1, f"{src.name}: estrazione testo")
    html = _office_to_html(src)

    if dst.suffix.lower() in (".md", ".markdown"):
        try:
            from markdownify import markdownify

            content = markdownify(html)
        except ImportError:
            content = BeautifulSoup(html, "html.parser").get_text("\n", strip=True)
    else:
        content = BeautifulSoup(html, "html.parser").get_text("\n", strip=True)

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(content, encoding="utf-8")
    ctx.progress(1, 1, f"{src.name}: completato")
    return [dst]


def spreadsheet_to_csv(
    src: Path,
    dst: Path,
    options: dict[str, Any] | None = None,
    ctx: JobContext = NULL_CONTEXT,
) -> list[Path]:
    """Export the first worksheet of a spreadsheet to CSV."""
    import csv

    import openpyxl

    src, dst = Path(src), Path(dst)
    ctx.progress(0, 1, f"{src.name}: lettura foglio")

    workbook = openpyxl.load_workbook(src, data_only=True, read_only=True)
    try:
        sheet = workbook.worksheets[0]
        dst.parent.mkdir(parents=True, exist_ok=True)
        with dst.open("w", encoding="utf-8-sig", newline="") as handle:
            writer = csv.writer(handle)
            for row in sheet.iter_rows(values_only=True):
                ctx.check_cancelled()
                writer.writerow(["" if cell is None else cell for cell in row])
    finally:
        workbook.close()

    ctx.progress(1, 1, f"{src.name}: completato")
    return [dst]
